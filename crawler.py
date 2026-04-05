"""
crawler.py

フォルダ内のファイルを読み込み、
各ファイルからテキストを抽出する。

対応ファイル:
- txt / md
- pdf（fitz使用）
- docx
- xls / xlsx
- pptx
"""

import re
from pathlib import Path
from datetime import datetime
from typing import List, Optional

import fitz                         # PDFを読み込むライブラリ
from docx import Document           # Wordファイルを読み込むライブラリ
from openpyxl import load_workbook  # Excel（xlsx）を読み込むライブラリ
import xlrd                         # Excel（xls）を読み込むライブラリ
from pptx import Presentation       # PowerPoint（pptx）を読み込むライブラリ
from ocr import ocr_pdf             # ocr.pyからocr_pdfを読み込むライブラリ


def normalize_text(text: str) -> str:
    """改行や空白を整理して読みやすくする"""
    if not text:
        return ""

    text = re.sub(r"\r\n|\r", "\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def extract_text_from_txt(file_path: Path) -> str:
    """txt / md ファイルを読み込む"""
    encodings = ["utf-8", "utf-8-sig", "cp932", "shift_jis", "latin-1"]

    for enc in encodings:
        try:
            return file_path.read_text(encoding=enc, errors="ignore")
        except Exception:
            continue

    return ""


def extract_text_from_pdf(file_path: Path) -> str:
    """PDF から本文を抽出する"""
    doc_pdf = fitz.open(file_path)
    texts = []

    for page in doc_pdf:
        try:
            page_text = page.get_text() or ""
            texts.append(page_text)
        except Exception:
            continue

    pdf_full_text = "\n".join(texts).strip()

    if len(pdf_full_text) >= 20:
        return pdf_full_text
    else:
        try:
            return ocr_pdf(file_path)
        except Exception as e:
            print(f"OCR処理中にエラーが発生しました: {e}")
            return ""


def extract_text_from_docx(file_path: Path) -> str:
    """Word(docx) から本文を抽出する"""
    doc = Document(file_path)
    texts = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            texts.append(text)

    return "\n".join(texts)


def extract_text_from_xlsx(file_path: Path) -> str:
    """Excel(xlsx) からセル内容を抽出する"""
    wb = load_workbook(file_path, data_only=True)
    texts = []

    for sheet in wb.worksheets:
        texts.append(f"[Sheet] {sheet.title}")

        for row in sheet.iter_rows(values_only=True):
            row_values = [
                str(cell).strip()
                for cell in row
                if cell is not None and str(cell).strip()
            ]
            if row_values:
                texts.append(" | ".join(row_values))

    return "\n".join(texts)


def extract_text_from_xls(file_path: Path) -> str:
    """Excel(xls) からセル内容を抽出する"""
    wb = xlrd.open_workbook(file_path)
    texts = []

    for sheet in wb.sheets():
        texts.append(f"[Sheet] {sheet.name}")

        for row_idx in range(sheet.nrows):
            row = sheet.row_values(row_idx)
            row_values = [str(cell).strip() for cell in row if str(cell).strip()]
            if row_values:
                texts.append(" | ".join(row_values))

    return "\n".join(texts)


def extract_text_from_pptx(file_path: Path) -> str:
    """PowerPoint(pptx) からスライド内テキストを抽出する"""
    prs = Presentation(file_path)
    texts = []

    for i, slide in enumerate(prs.slides, start=1):
        texts.append(f"[Slide] {i}")

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    texts.append(text)

    return "\n".join(texts)


def clean_author_value(value) -> str:
    """author用の値を安全に文字列化する"""
    if value is None:
        return ""
    value = str(value).strip()
    if not value:
        return ""
    return value


def extract_author(file_path: Path, suffix: str) -> str:
    """ファイルの作成者 / 最終更新者をできるだけ取得する"""
    try:
        if suffix == ".docx":
            doc = Document(file_path)
            cp = doc.core_properties

            last_modified_by = clean_author_value(getattr(cp, "last_modified_by", None))
            if last_modified_by:
                return last_modified_by

            author = clean_author_value(getattr(cp, "author", None))
            if author:
                return author

            return "不明"

        elif suffix == ".xlsx":
            wb = load_workbook(file_path)
            props = wb.properties

            # openpyxlでは lastModifiedBy のことが多い
            last_modified_by = clean_author_value(getattr(props, "lastModifiedBy", None))
            if last_modified_by:
                return last_modified_by

            # 念のため snake_case も確認
            last_modified_by = clean_author_value(getattr(props, "last_modified_by", None))
            if last_modified_by:
                return last_modified_by

            creator = clean_author_value(getattr(props, "creator", None))
            if creator:
                return creator

            return "不明"

        elif suffix == ".xls":
            wb = xlrd.open_workbook(file_path)
            user_name = clean_author_value(getattr(wb, "user_name", None))
            if user_name:
                return user_name
            return "不明"

        elif suffix == ".pptx":
            prs = Presentation(file_path)
            cp = prs.core_properties

            last_modified_by = clean_author_value(getattr(cp, "last_modified_by", None))
            if last_modified_by:
                return last_modified_by

            author = clean_author_value(getattr(cp, "author", None))
            if author:
                return author

            return "不明"

        elif suffix == ".pdf":
            doc_pdf = fitz.open(file_path)
            meta = doc_pdf.metadata or {}

            author = clean_author_value(meta.get("author"))
            if author:
                return author

            creator = clean_author_value(meta.get("creator"))
            if creator:
                return creator

            producer = clean_author_value(meta.get("producer"))
            if producer:
                return producer

            return "不明"

        return "不明"

    except Exception:
        return "不明"


def infer_description(text: str, max_len: int = 200) -> str:
    """本文先頭から説明文を作る"""
    text = normalize_text(text)
    return text[:max_len]


def extract_file_content(file_path: Path) -> dict:
    """
    1ファイル分を処理し、
    app.py / database.py が使いやすい形式で返す
    """
    updated_at = datetime.fromtimestamp(file_path.stat().st_mtime)
    suffix = file_path.suffix.lower()

    # --- author取得 ---
    author = extract_author(file_path, suffix)

    # --- 本文抽出 ---
    try:
        if suffix in [".txt", ".md"]:
            raw_text = extract_text_from_txt(file_path)
        elif suffix == ".pdf":
            raw_text = extract_text_from_pdf(file_path)
        elif suffix == ".docx":
            raw_text = extract_text_from_docx(file_path)
        elif suffix == ".xlsx":
            raw_text = extract_text_from_xlsx(file_path)
        elif suffix == ".xls":
            raw_text = extract_text_from_xls(file_path)
        elif suffix == ".pptx":
            raw_text = extract_text_from_pptx(file_path)
        else:
            return {
                "url": str(file_path.resolve()),
                "title": file_path.stem,
                "file_name": file_path.name,
                "file_type": suffix.lstrip("."),
                "description": "",
                "keywords": "",
                "full_text": "",
                "author": author,
                "category": file_path.parent.name,
                "word_count": 0,
                "crawled_at": datetime.now().isoformat(),
                "created_at": updated_at.isoformat(),
                "crawl_status": "failed",
                "error_message": f"未対応形式: {suffix}",
            }

        full_text = normalize_text(raw_text)

        return {
            "url": str(file_path.resolve()),
            "title": file_path.stem,
            "file_name": file_path.name,
            "file_type": suffix.lstrip("."),
            "description": infer_description(full_text),
            "keywords": "",
            "full_text": full_text,
            "author": author,
            "category": file_path.parent.name,
            "word_count": len(full_text.split()),
            "crawled_at": datetime.now().isoformat(),
            "created_at": updated_at.isoformat(),
            "crawl_status": "success",
            "error_message": "",
        }

    except Exception as e:
        return {
            "url": str(file_path.resolve()),
            "title": file_path.stem,
            "file_name": file_path.name,
            "file_type": suffix.lstrip("."),
            "description": "",
            "keywords": "",
            "full_text": "",
            "author": author,
            "category": file_path.parent.name,
            "word_count": 0,
            "crawled_at": datetime.now().isoformat(),
            "created_at": updated_at.isoformat(),
            "crawl_status": "failed",
            "error_message": str(e),
        }


def list_files(
    folder_path: str,
    recursive: bool = True,
    allowed_extensions: Optional[List[str]] = None,
    max_files: int = 200
) -> List[Path]:
    """フォルダ内の対象ファイル一覧を取得する"""
    base = Path(folder_path)

    if allowed_extensions is None:
        allowed_extensions = [".txt", ".md", ".pdf", ".docx", ".xls", ".xlsx", ".pptx"]

    allowed_set = {ext.lower() for ext in allowed_extensions}
    iterator = base.rglob("*") if recursive else base.glob("*")

    files = []
    for p in iterator:
        if p.is_file() and p.suffix.lower() in allowed_set:
            files.append(p)
            if len(files) >= max_files:
                break

    return files


def import_folder(
    folder_path: str,
    recursive: bool = True,
    allowed_extensions: Optional[List[str]] = None,
    max_files: int = 200
) -> List[dict]:
    """フォルダ内ファイルパスをリストにまとめて抽出する"""
    files = list_files(
        folder_path=folder_path,
        recursive=recursive,
        allowed_extensions=allowed_extensions,
        max_files=max_files
    )

    results = []
    for file_path in files:
        results.append(extract_file_content(file_path))

    return results